import os
import sys
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("python-pptx is not installed. Attempting to install it now...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
        print("python-pptx successfully installed!")
    except Exception as e:
        print(f"Failed to install python-pptx: {e}")
        print("Please run: pip install python-pptx")
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Palette definition
    c_bg_main = RGBColor(11, 15, 25)       # Slate dark background (#0b0f19)
    c_bg_card = RGBColor(30, 41, 59)       # Card background (#1e293b)
    c_border = RGBColor(51, 65, 85)        # Border/line color (#334155)
    c_text_main = RGBColor(248, 250, 252)  # White text (#f8fafc)
    c_text_muted = RGBColor(148, 163, 184) # Muted text (#94a3b8)
    
    c_cyan = RGBColor(6, 182, 212)         # Cyan (#06b6d4)
    c_cyan_light = RGBColor(34, 211, 238)   # Cyan light (#22d3ee)
    c_emerald = RGBColor(16, 185, 129)     # Emerald (#10b981)
    c_rose = RGBColor(244, 63, 94)         # Rose (#f43f5e)
    c_amber = RGBColor(245, 158, 11)       # Amber (#f59e0b)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    def apply_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_bg_main

    def add_slide_header(slide, title_text, slide_num):
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(10), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = c_cyan
        
        # Slide Number Tag
        num_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.45), Inches(1.08), Inches(0.4))
        tf_num = num_box.text_frame
        tf_num.word_wrap = True
        tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
        p_num = tf_num.paragraphs[0]
        p_num.alignment = PP_ALIGN.RIGHT
        p_num.text = f"{slide_num:02d} / 12"
        p_num.font.name = 'Courier New'
        p_num.font.size = Pt(14)
        p_num.font.bold = True
        p_num.font.color.rgb = c_cyan
        
        # Divider Line
        connector = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.2), Inches(11.83), Inches(0.02)
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = c_border
        connector.line.color.rgb = c_border

    def add_bullet_list(slide, items, left, top, width, height, font_size=16):
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        for idx, item in enumerate(items):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = "•  " + item[0]
            p.font.name = 'Arial'
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = c_text_main
            p.space_after = Pt(4)
            
            p_desc = tf.add_paragraph()
            p_desc.text = "    " + item[1]
            p_desc.font.name = 'Arial'
            p_desc.font.size = Pt(font_size - 2)
            p_desc.font.color.rgb = c_text_muted
            p_desc.space_after = Pt(14)
            
        return tx_box

    def add_card(slide, title, text, left, top, width, height, highlight_color):
        # Card background shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = c_bg_card
        card.line.color.rgb = c_border
        card.line.width = Pt(1)
        
        # Left accent line
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.1), Inches(0.08), height - Inches(0.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = highlight_color
        accent.line.fill.background()
        
        # Text Frame
        tx_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.45), height - Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = highlight_color
        p.space_after = Pt(8)
        
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.name = 'Arial'
        p2.font.size = Pt(13)
        p2.font.color.rgb = c_text_muted
        
        return card

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide1)
    
    # Large Decorative Glow shape in background (subtle circle)
    glow = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(12, 35, 55)
    glow.line.fill.background()
    
    # Title text
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.83), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "IoT Network Traffic Generator"
    p.font.name = 'Arial'
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = c_text_main
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "A Modular, Multi-Protocol Simulation Framework"
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.color.rgb = c_cyan_light
    
    # Info Grid Cards
    cards_data = [
        ("Submitted By", "Amjad Khan\nCS & Engineering Student"),
        ("Supervised By", "Dr. Bhale Pradeepkumar Gajendra\nIIITDM Kancheepuram"),
        ("Program", "AICTE Smarthan Internship\n2026-2027"),
        ("Host Institution", "IIITDM Kancheepuram\nChennai, India")
    ]
    
    for idx, (label, val) in enumerate(cards_data):
        row = idx // 2
        col = idx % 2
        
        x = Inches(0.75) + col * Inches(6.0)
        y = Inches(3.8) + row * Inches(1.5)
        w = Inches(5.8)
        h = Inches(1.2)
        
        card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = c_bg_card
        card.line.color.rgb = c_border
        
        tx = slide1.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), h - Inches(0.3))
        tf_c = tx.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = 0
        
        p_lbl = tf_c.paragraphs[0]
        p_lbl.text = label.upper()
        p_lbl.font.name = 'Arial'
        p_lbl.font.size = Pt(11)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = c_cyan
        p_lbl.space_after = Pt(4)
        
        p_val = tf_c.add_paragraph()
        p_val.text = val
        p_val.font.name = 'Arial'
        p_val.font.size = Pt(14)
        p_val.font.bold = True
        p_val.font.color.rgb = c_text_main
        
    # ==========================================
    # SLIDE 2: Agenda
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide2)
    add_slide_header(slide2, "Agenda & Outline", 2)
    
    col1_items = [
        ("01", "Introduction (Background & Motivation)"),
        ("02", "Problem Statement & Objectives"),
        ("03", "Related Work & Literature Review"),
        ("04", "Proposed Work & Implementation Details"),
    ]
    col2_items = [
        ("05", "System Design & Architecture Diagram"),
        ("06", "Testing, Results & Performance Evaluation"),
        ("07", "Discussion & Load Impact Analysis"),
        ("08", "Conclusion, Future Scope & References"),
    ]
    
    def add_timeline_col(slide, items, left, top_start):
        for idx, (num, text) in enumerate(items):
            y = top_start + idx * Inches(1.1)
            # Draw number box
            num_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, Inches(0.75), Inches(0.75))
            num_shape.fill.solid()
            num_shape.fill.fore_color.rgb = c_bg_card
            num_shape.line.color.rgb = c_cyan
            
            tf_num = num_shape.text_frame
            tf_num.word_wrap = True
            p_n = tf_num.paragraphs[0]
            p_n.alignment = PP_ALIGN.CENTER
            p_n.text = num
            p_n.font.name = 'Courier New'
            p_n.font.size = Pt(20)
            p_n.font.bold = True
            p_n.font.color.rgb = c_cyan
            
            # Draw text box
            tx_shape = slide.shapes.add_textbox(left + Inches(0.95), y + Inches(0.12), Inches(4.85), Inches(0.55))
            tf_tx = tx_shape.text_frame
            tf_tx.word_wrap = True
            tf_tx.margin_left = tf_tx.margin_right = tf_tx.margin_top = tf_tx.margin_bottom = 0
            p_t = tf_tx.paragraphs[0]
            p_t.text = text
            p_t.font.name = 'Arial'
            p_t.font.size = Pt(16)
            p_t.font.bold = True
            p_t.font.color.rgb = c_text_main
            
    add_timeline_col(slide2, col1_items, Inches(0.75), Inches(2.0))
    add_timeline_col(slide2, col2_items, Inches(6.8), Inches(2.0))

    # ==========================================
    # SLIDE 3: Introduction
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide3)
    add_slide_header(slide3, "Introduction", 3)
    
    # Left column: background info
    left_tx = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_l = left_tx.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "Background & Motivation"
    p_lh.font.name = 'Arial'
    p_lh.font.size = Pt(22)
    p_lh.font.bold = True
    p_lh.font.color.rgb = c_cyan_light
    p_lh.space_after = Pt(12)
    
    p_lbody = tf_l.add_paragraph()
    p_lbody.text = "With the exponential rise of IoT devices across smart homes, industrial complexes, and medical infrastructures, networks face massive scaling demands. Simulating network traffic enables testing infrastructures before actual deployment."
    p_lbody.font.name = 'Arial'
    p_lbody.font.size = Pt(16)
    p_lbody.font.color.rgb = c_text_muted
    p_lbody.space_after = Pt(24)
    
    # Draw Highlight Tip Box
    tip_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.4), Inches(5.5), Inches(1.8))
    tip_card.fill.solid()
    tip_card.fill.fore_color.rgb = RGBColor(12, 35, 55)
    tip_card.line.color.rgb = c_cyan
    
    tf_tip = tip_card.text_frame
    tf_tip.word_wrap = True
    tf_tip.margin_left = tf_tip.margin_right = tf_tip.margin_top = tf_tip.margin_bottom = Inches(0.15)
    p_t1 = tf_tip.paragraphs[0]
    p_t1.text = "KEY RESEARCH FACTOR"
    p_t1.font.name = 'Arial'
    p_t1.font.size = Pt(12)
    p_t1.font.bold = True
    p_t1.font.color.rgb = c_cyan
    p_t1.space_after = Pt(6)
    p_t2 = tf_tip.add_paragraph()
    p_t2.text = "Software-defined network simulation offers a highly scalable, zero-dependency, cost-effective alternative to deploying physical hardware loops."
    p_t2.font.name = 'Arial'
    p_t2.font.size = Pt(13)
    p_t2.font.color.rgb = c_text_main
    
    # Right column: bullet lists
    right_items = [
        ("Heterogeneous Protocols", "IoT simulation demands supporting diverse application layer protocols: UDP (low latency), TCP (reliable streams), HTTP (REST interfaces), and MQTT (publish-subscribe telemetry)."),
        ("Variable Traffic Dynamics", "Device traffic profiles fluctuate between standard telemetry routines, critical event bursts, and malware-driven floods.")
    ]
    add_bullet_list(slide3, right_items, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 4: Problem Statement
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide4)
    add_slide_header(slide4, "Problem Statement", 4)
    
    # Cards layout
    add_card(slide4, "Cost & Logistics", 
             "Physical verification of setups using hundreds of distinct IoT devices is cost-prohibitive and difficult to automate or configure.", 
             Inches(0.75), Inches(1.8), Inches(3.7), Inches(3.2), c_rose)
             
    add_card(slide4, "Traffic Diversity", 
             "Standard simulators utilize homogeneous, static patterns, failing to model compound multi-protocol payloads or sudden state changes.", 
             Inches(4.8), Inches(1.8), Inches(3.7), Inches(3.2), c_amber)
             
    add_card(slide4, "Security Stress Testing", 
             "IoT units are vulnerable to botnets like Mirai. Security researchers require dynamic DDoS workloads to verify firewall rules.", 
             Inches(8.85), Inches(1.8), Inches(3.7), Inches(3.2), c_rose)
             
    # Bottom warning bar
    warning_bar = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.3), Inches(11.83), Inches(1.2))
    warning_bar.fill.solid()
    warning_bar.fill.fore_color.rgb = RGBColor(40, 15, 25)
    warning_bar.line.color.rgb = c_rose
    
    tf_w = warning_bar.text_frame
    tf_w.word_wrap = True
    tf_w.margin_left = tf_w.margin_right = tf_w.margin_top = tf_w.margin_bottom = Inches(0.15)
    
    p_w1 = tf_w.paragraphs[0]
    p_w1.text = "RESEARCH GAP"
    p_w1.font.name = 'Arial'
    p_w1.font.size = Pt(12)
    p_w1.font.bold = True
    p_w1.font.color.rgb = c_rose
    p_w1.space_after = Pt(4)
    
    p_w2 = tf_w.add_paragraph()
    p_w2.text = "There is a critical lack of zero-dependency, Python-only simulators capable of shifting device profiles dynamically from normal periodic updates to burst alerts and botnet attack loads during runtime."
    p_w2.font.name = 'Arial'
    p_w2.font.size = Pt(14)
    p_w2.font.color.rgb = c_text_main

    # ==========================================
    # SLIDE 5: Objectives
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide5)
    add_slide_header(slide5, "Internship Objectives", 5)
    
    left_obj = [
        ("Protocol Exploration", "Deep dive into framing constraints, packet structures, and transfer protocols (UDP, TCP, HTTP, MQTT)."),
        ("Concurrency Simulation", "Implement a robust multi-threaded client model in Python that manages concurrent device loops."),
        ("Multi-Profile Traffic Modeling", "Formulate realistic schedules for Normal, Burst, and DDoS Attack traffic types.")
    ]
    right_obj = [
        ("Low-Level socket implementation", "Code raw byte-level MQTT CONNECT and PUBLISH packets to bypass third-party library requirements."),
        ("Aggregation and Statistics Tracker", "Build a receiver server that aggregates socket buffers, tracks packets, and prints network performance metrics.")
    ]
    
    add_bullet_list(slide5, left_obj, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    add_bullet_list(slide5, right_obj, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 6: Related Work
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide6)
    add_slide_header(slide6, "Related Work & Literature Review", 6)
    
    # Left Box
    box_l = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_bl = box_l.text_frame
    tf_bl.word_wrap = True
    tf_bl.margin_left = tf_bl.margin_right = tf_bl.margin_top = tf_bl.margin_bottom = 0
    
    p_lh6 = tf_bl.paragraphs[0]
    p_lh6.text = "Standards & RFC Documentation"
    p_lh6.font.name = 'Arial'
    p_lh6.font.size = Pt(20)
    p_lh6.font.bold = True
    p_lh6.font.color.rgb = c_cyan_light
    p_lh6.space_after = Pt(14)
    
    docs = [
        ("RFC 768 (UDP)", "Connectionless structure, packet headers, MTU limits, and transmission mechanics."),
        ("RFC 2616 (HTTP/1.1)", "REST methods, headers, payload framing, and TCP socket persistence."),
        ("MQTT v3.1.1 (OASIS)", "Binary wire format protocol, publish/subscribe schemas, control bytes, and topic parsing.")
    ]
    for d_title, d_desc in docs:
        p_t = tf_bl.add_paragraph()
        p_t.text = "•  " + d_title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = c_text_main
        
        p_d = tf_bl.add_paragraph()
        p_d.text = "    " + d_desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = c_text_muted
        p_d.space_after = Pt(10)
        
    # Right Box
    box_r = slide6.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    tf_br = box_r.text_frame
    tf_br.word_wrap = True
    tf_br.margin_left = tf_br.margin_right = tf_br.margin_top = tf_br.margin_bottom = 0
    
    p_rh6 = tf_br.paragraphs[0]
    p_rh6.text = "Literature & Simulators"
    p_rh6.font.name = 'Arial'
    p_rh6.font.size = Pt(20)
    p_rh6.font.bold = True
    p_rh6.font.color.rgb = c_cyan_light
    p_rh6.space_after = Pt(14)
    
    p_lit1 = tf_br.add_paragraph()
    p_lit1.text = "Botnet Traffic Emulation"
    p_lit1.font.name = 'Arial'
    p_lit1.font.size = Pt(16)
    p_lit1.font.bold = True
    p_lit1.font.color.rgb = c_text_main
    p_lit1.space_after = Pt(4)
    
    p_lit2 = tf_br.add_paragraph()
    p_lit2.text = "Research on the Mirai Botnet (IEEE S&P, 2017) details how vulnerable IoT units form coordinate swarms. This informs our HTTP/TCP flood logic."
    p_lit2.font.name = 'Arial'
    p_lit2.font.size = Pt(14)
    p_lit2.font.color.rgb = c_text_muted
    p_lit2.space_after = Pt(20)
    
    p_lit3 = tf_br.add_paragraph()
    p_lit3.text = "Simulators Comparison"
    p_lit3.font.name = 'Arial'
    p_lit3.font.size = Pt(16)
    p_lit3.font.bold = True
    p_lit3.font.color.rgb = c_text_main
    p_lit3.space_after = Pt(4)
    
    p_lit4 = tf_br.add_paragraph()
    p_lit4.text = "Standard tools like NS3, Mininet, or PcapPlay require hypervisors or custom operating environments, lacking lightweight cross-platform scripts."
    p_lit4.font.name = 'Arial'
    p_lit4.font.size = Pt(14)
    p_lit4.font.color.rgb = c_text_muted

    # ==========================================
    # SLIDE 7: Proposed Work
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide7)
    add_slide_header(slide7, "Proposed Work & Implementation Details", 7)
    
    # Left column: Text details
    left_tx7 = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_l7 = left_tx7.text_frame
    tf_l7.word_wrap = True
    tf_l7.margin_left = tf_l7.margin_right = tf_l7.margin_top = tf_l7.margin_bottom = 0
    
    p_lh7 = tf_l7.paragraphs[0]
    p_lh7.text = "Python-Based Client-Server Architecture"
    p_lh7.font.name = 'Arial'
    p_lh7.font.size = Pt(20)
    p_lh7.font.bold = True
    p_lh7.font.color.rgb = c_cyan_light
    p_lh7.space_after = Pt(14)
    
    bullets = [
        ("Multi-threading Model", "Each device runs in a separate thread monitoring real-time configuration overrides."),
        ("Dynamic JSON Config Tracker", "A checker loop re-reads `config.json` dynamic parameters without system reboots."),
        ("Raw Byte MQTT Packing", "Custom socket buffer methods create MQTT payloads directly at byte level without importing complex packages.")
    ]
    for b_lbl, b_desc in bullets:
        p_t = tf_l7.add_paragraph()
        p_t.text = "•  " + b_lbl
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = c_text_main
        
        p_d = tf_l7.add_paragraph()
        p_d.text = "    " + b_desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = c_text_muted
        p_d.space_after = Pt(10)
        
    # Right column: Styled Code Box
    code_card = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    code_card.fill.solid()
    code_card.fill.fore_color.rgb = RGBColor(8, 12, 20)
    code_card.line.color.rgb = c_border
    
    tf_code = code_card.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = tf_code.margin_right = tf_code.margin_top = tf_code.margin_bottom = Inches(0.2)
    
    p_cd0 = tf_code.paragraphs[0]
    p_cd0.text = "# Raw MQTT packet assembler"
    p_cd0.font.name = 'Courier New'
    p_cd0.font.size = Pt(12)
    p_cd0.font.color.rgb = RGBColor(100, 116, 139) # comment color
    p_cd0.space_after = Pt(4)
    
    code_lines = [
        ("def make_mqtt_publish_packet(topic, message):", RGBColor(244, 63, 94)),
        ("    topic_bytes = topic.encode('utf-8')", RGBColor(226, 232, 240)),
        ("    topic_len = len(topic_bytes).to_bytes(2, 'big')", RGBColor(226, 232, 240)),
        ("    msg_bytes = message.encode('utf-8')", RGBColor(226, 232, 240)),
        ("    payload = topic_len + topic_bytes + msg_bytes", RGBColor(226, 232, 240)),
        ("    rem_len = encode_remaining_length(len(payload))", RGBColor(226, 232, 240)),
        ("    # 0x30 = PUBLISH Control Type", RGBColor(100, 116, 139)),
        ("    return bytes([0x30]) + rem_len + payload", RGBColor(244, 63, 94))
    ]
    for text, color in code_lines:
        p_line = tf_code.add_paragraph()
        p_line.text = text
        p_line.font.name = 'Courier New'
        p_line.font.size = Pt(12)
        p_line.font.color.rgb = color
        p_line.space_after = Pt(2)

    # ==========================================
    # SLIDE 8: System Architecture
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide8)
    add_slide_header(slide8, "System Architecture", 8)
    
    # Check if the PNG diagram exists
    img_path = 'mermaid-diagram-2026-07-07-142113.png'
    if os.path.exists(img_path):
        # Insert image in the center of the slide
        slide8.shapes.add_picture(img_path, Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.8))
    else:
        # Fallback text box if image doesn't exist
        card_fallback = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(2.2), Inches(9.33), Inches(3.5))
        card_fallback.fill.solid()
        card_fallback.fill.fore_color.rgb = c_bg_card
        card_fallback.line.color.rgb = c_border
        
        tf_fb = card_fallback.text_frame
        tf_fb.word_wrap = True
        p_fb = tf_fb.paragraphs[0]
        p_fb.alignment = PP_ALIGN.CENTER
        p_fb.text = "[ Architecture Diagram Placeholder ]\n\n(Ensure mermaid-diagram-2026-07-07-142113.png exists in this folder for the script to load it automatically.)"
        p_fb.font.name = 'Arial'
        p_fb.font.size = Pt(16)
        p_fb.font.color.rgb = c_text_muted

    # ==========================================
    # SLIDE 9: Testing & Results
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide9)
    add_slide_header(slide9, "Testing & Results", 9)
    
    # Table layout
    rows = 4
    cols = 5
    left = Inches(0.75)
    top = Inches(2.2)
    width = Inches(11.83)
    height = Inches(3.5)
    
    table_shape = slide9.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Define columns width
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.2)
    table.columns[4].width = Inches(2.73)
    
    headers = ["Traffic Profile", "Packets Received", "Total Bytes", "Avg Byte Rate", "Server Status"]
    data = [
        ["Normal", "108", "18,360", "~306 B/s", "Stable (Zero Loss)"],
        ["Burst", "432", "75,168", "~1.25 KB/s", "Stable (Low Latency)"],
        ["Attack (DDoS)", "14,210", "4,263,000", "~71.05 KB/s", "High Load (Delayed)"]
    ]
    
    # Style Table Header
    for c_idx, text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Arial'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = c_cyan
        
    # Style Table Body
    colors_accent = [c_emerald, c_amber, c_rose]
    for r_idx, row_data in enumerate(data):
        cell_accent = colors_accent[r_idx]
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(18, 24, 38)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Arial'
            p.font.size = Pt(13)
            p.font.bold = (c_idx == 0) # Bold profile name
            p.font.color.rgb = cell_accent if (c_idx == 0 or c_idx == 4) else c_text_main

    # ==========================================
    # SLIDE 10: Discussion
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide10)
    add_slide_header(slide10, "Discussion", 10)
    
    add_card(slide10, "Socket Level Optimization", 
             "Bypassing third-party modules enabled the Python interpreter to assemble network headers with close-to-zero memory and process footprint overheads.", 
             Inches(0.75), Inches(1.8), Inches(5.6), Inches(2.8), c_emerald)
             
    add_card(slide10, "Dynamic Adaptability", 
             "The local config tracking logic reacted in milliseconds. Transitioning nodes to Burst triggered immediate message frequency multipliers.", 
             Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.8), c_amber)
             
    # Stress Testing Highlight Card
    stress_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.9), Inches(11.83), Inches(1.6))
    stress_card.fill.solid()
    stress_card.fill.fore_color.rgb = RGBColor(40, 15, 25)
    stress_card.line.color.rgb = c_rose
    
    tf_sc = stress_card.text_frame
    tf_sc.word_wrap = True
    tf_sc.margin_left = tf_sc.margin_right = tf_sc.margin_top = tf_sc.margin_bottom = Inches(0.15)
    
    p_sc1 = tf_sc.paragraphs[0]
    p_sc1.text = "DDoS EVALUATION & ATTACK THRESHOLDS"
    p_sc1.font.name = 'Arial'
    p_sc1.font.size = Pt(14)
    p_sc1.font.bold = True
    p_sc1.font.color.rgb = c_rose
    p_sc1.space_after = Pt(4)
    
    p_sc2 = tf_sc.add_paragraph()
    p_sc2.text = "Under the Attack Profile (14,210 packets sent over 60s), the receiver server threads encountered message queue delays. This validated the simulator's ability to trigger network stress states, providing a cost-effective, zero-dependency environment for testing system resilience and log analyzer rules."
    p_sc2.font.name = 'Arial'
    p_sc2.font.size = Pt(13)
    p_sc2.font.color.rgb = c_text_main

    # ==========================================
    # SLIDE 11: Conclusion & Future Scope
    # ==========================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide11)
    add_slide_header(slide11, "Conclusion & Future Scope", 11)
    
    left_c = [
        ("Multi-Protocol Simulation Engine", "Successfully built a robust traffic generator capable of HTTP, TCP, UDP, and MQTT emulation."),
        ("Active IDS/IPS Firewall Engine", "Implemented real-time DDoS rate limiting (>12 pkt/s) and signature-based packet dropping."),
        ("Pure-Python PCAP Logging", "Engineered byte-level PCAP header construction for Wireshark inspection without external C dependencies.")
    ]
    right_c = [
        ("Machine Learning IDS", "Replace the rule-based rate limiter with an ML model (e.g., Isolation Forest) for adaptive anomaly detection."),
        ("Real Device Integration", "Extend the generator to communicate with actual hardware IoT sensors (Raspberry Pi, ESP32)."),
        ("Cloud Deployment & IPv6", "Deploy the receiver on a cloud VM and upgrade the PCAP logger to construct IPv6 headers.")
    ]
    
    add_bullet_list(slide11, left_c, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5), font_size=15)
    add_bullet_list(slide11, right_c, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5), font_size=15)

    # ==========================================
    # SLIDE 12: References
    # ==========================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide12)
    add_slide_header(slide12, "References", 12)
    
    ref_items = [
        ("Python Foundation", "Python 3 Docs: socket, threading, struct, http.server modules."),
        ("OASIS Standard", "MQTT Version 3.1.1 Specification: Publish/Subscribe Messaging Protocol for IoT."),
        ("IETF (RFC 791, 793, 768)", "Internet Protocol (IPv4), TCP, and UDP Specifications."),
        ("Wireshark Foundation", "Libpcap File Format and PCAP Header Encoding Standards."),
        ("V. Chandola et al.", "Anomaly Detection: A Survey, ACM Computing Surveys, 2009.")
    ]
    
    tx_box = slide12.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    for idx, (auth, pub) in enumerate(ref_items):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = f"[{idx+1}]  "
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = c_cyan
        
        # Add Author
        run_auth = p.add_run()
        run_auth.text = f"{auth}, "
        run_auth.font.bold = True
        run_auth.font.color.rgb = c_text_main
        
        # Add Pub info
        run_pub = p.add_run()
        run_pub.text = pub
        run_pub.font.bold = False
        run_pub.font.color.rgb = c_text_muted
        
        p.space_after = Pt(14)
        
    # Save the presentation
    output_filename = "presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'!")

if __name__ == "__main__":
    create_presentation()
